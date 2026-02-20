import sys
import subprocess

# Check if python-pptx is installed
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("Error: 'python-pptx' library is not installed.")
    print("Please install it by running: pip install python-pptx")
    sys.exit(1)

def create_presentation():
    prs = Presentation()

    # Define content
    slides_content = [
        {
            "title": "Document RAG System",
            "subtitle": "Interactive Chat with your Documents",
            "type": "title"
        },
        {
            "title": "Agenda",
            "content": [
                "1. Introduction",
                "2. Problem Statement",
                "3. Solution Overview",
                "4. System Architecture",
                "5. Tech Stack",
                "6. Key Features",
                "7. Implementation Details",
                "8. Future Roadmap",
                "9. Q&A"
            ],
            "type": "content"
        },
        {
            "title": "Introduction",
            "content": [
                "What is RAG? (Retrieval-Augmented Generation)",
                "- Combines LLMs with proprietary data.",
                "- Ensures answers are grounded in facts.",
                "",
                "Why it matters?",
                "- Enables AI to answer questions about YOUR documents.",
                "- Overcomes the knowledge cutoff of standard models."
            ],
            "type": "content"
        },
        {
            "title": "Problem Statement",
            "content": [
                "Limitations of Standard LLMs:",
                "- Lack knowledge of private/specific data.",
                "- Prone to hallucinations.",
                "",
                "Manual Document Search:",
                "- Time-consuming to read through PDFs.",
                "- Keyword search (Ctrl+F) misses semantic meaning."
            ],
            "type": "content"
        },
        {
            "title": "Solution Overview",
            "content": [
                "Document Ingestion: Upload PDF, DOCX, TXT.",
                "Semantic Search: Understands query meaning.",
                "Generative Answers: Logic + Data -> Answer.",
                "Result: Instant, cited, and accurate responses."
            ],
            "type": "content"
        },
        {
            "title": "System Architecture",
            "content": [
                "1. Ingestion Pipeline:",
                "   File -> Text -> Chunks -> Embeddings -> ChromaDB",
                "",
                "2. Query Pipeline:",
                "   Question -> Rewrite -> Retrieval -> Generation -> Answer"
            ],
            "type": "content"
        },
        {
            "title": "Tech Stack",
            "content": [
                "Backend: Python, FastAPI",
                "Frontend: HTML/CSS/JS (Jinja2)",
                "Vector DB: ChromaDB",
                "LLM: Google Gemini (via API)",
                "Embeddings: SentenceTransformers (all-MiniLM-L6-v2)"
            ],
            "type": "content"
        },
        {
            "title": "Key Features",
            "content": [
                "📂 Multi-format Support (PDF, DOCX, TXT)",
                "💬 Context-Aware Chat History",
                "🔊 Text-to-Speech (TTS) Integration",
                "🔗 Source Citations (Hidden for cleaner UI)",
                "⚡ Background File Processing",
                "🗑️ Easy File & Chat Management"
            ],
            "type": "content"
        },
        {
            "title": "Implementation Details",
            "content": [
                "Chunking: 1000 chars with 200 overlap.",
                "Prompting: Strict context adherence to reduce hallucinations.",
                "Infrastructure: Local Setup."
            ],
            "type": "content"
        },
        {
            "title": "Future Roadmap",
            "content": [
                "🎙️ Voice Interface (Speech-to-Text)",
                "🔐 User Authentication",
                "📊 Excel/CSV Support",
                "🧠 Hybrid Search (Keyword + Semantic)"
            ],
            "type": "content"
        },
        {
            "title": "Q&A",
            "content": [
                "Open Floor for Questions",
                "Thank You!"
            ],
            "type": "content"
        }
    ]

    # Create slides
    for slide_data in slides_content:
        if slide_data["type"] == "title":
            slide_layout = prs.slide_layouts[0] # Title Slide
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = slide_data["title"]
            subtitle.text = slide_data["subtitle"]
        
        elif slide_data["type"] == "content":
            slide_layout = prs.slide_layouts[1] # Title and Content
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            body = slide.placeholders[1]
            title.text = slide_data["title"]
            
            tf = body.text_frame
            tf.word_wrap = True
            
            for line in slide_data["content"]:
                p = tf.add_paragraph()
                p.text = line
                if line.startswith("-"):
                    p.level = 1

    # Save
    output_file = "Document_RAG_Presentation.pptx"
    prs.save(output_file)
    print(f"Presentation saved to {output_file}")

if __name__ == "__main__":
    create_presentation()
